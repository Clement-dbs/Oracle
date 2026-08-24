"""
Configuration pytest partagée.

Les variables d'environnement nécessaires à app.core.config sont fixées
ICI, avant tout import de `app.*` par un module de test -- config.py fait
des int(os.getenv(...)) sans valeur par défaut : sans ces variables,
n'importe quel `from app... import ...` lève une TypeError au chargement.
Aucun service externe (Qdrant, MinIO, Redis, Ollama) n'est requis pour
lancer cette suite : les clients sont mockés dans chaque test qui en a
besoin, ces variables ne servent qu'à permettre l'import des modules.
"""

import datetime as _datetime
import os

# `datetime.UTC` n'existe qu'à partir de Python 3.11 (le code de production
# cible 3.11+, voir Dockerfile) -- ce shim permet de lancer la suite de tests
# aussi sous un interpréteur 3.10, sans changer le code applicatif pour
# autant (qui doit rester idiomatique pour sa cible réelle).
if not hasattr(_datetime, "UTC"):
    _datetime.UTC = _datetime.timezone.utc  # noqa: UP017 -- c'est le shim lui-même

os.environ["REDIS_URL"] = "redis://localhost:6379/0"
os.environ["QDRANT_URL"] = "http://localhost:6333"
os.environ["QDRANT_COLLECTION"] = "test_collection"
os.environ["EMBEDDINGS_MODEL"] = "BAAI/bge-m3"
os.environ["EMBEDDINGS_MODEL_LOCAL_PATH"] = "/tmp/test_embeddings_model"
os.environ["RERANKER_MODEL"] = "BAAI/bge-reranker-v2-m3"
os.environ["RERANKER_LOCAL_PATH"] = "/tmp/test_reranker_model"
os.environ["OLLAMA_HOST"] = "http://localhost:11434"
os.environ["OLLAMA_MODEL"] = "qwen2.5:7b"
os.environ["MINIO_ENDPOINT"] = "localhost:9000"
os.environ["MINIO_ROOT_USER"] = "test"
os.environ["MINIO_ROOT_PASSWORD"] = "test1234"
os.environ["MINIO_BUCKET"] = "test-bucket"

import sys  # noqa: E402
import types  # noqa: E402

# FlagEmbedding (BGE-M3) tire torch/CUDA -- bien trop lourd pour des tests
# unitaires et inutile ici : chaque test qui a besoin d'un modèle mocke
# get_model()/_model directement. On stub juste le module pour que
# `from FlagEmbedding import BGEM3FlagModel` (app/ingestion/embeddings.py)
# ne fasse pas échouer l'import si le paquet réel n'est pas installé.
if "FlagEmbedding" not in sys.modules:
    _fake_flagembedding = types.ModuleType("FlagEmbedding")
    _fake_flagembedding.BGEM3FlagModel = object
    sys.modules["FlagEmbedding"] = _fake_flagembedding

import fitz  # noqa: E402
import pytest  # noqa: E402


@pytest.fixture
def native_pdf_path(tmp_path):
    """PDF d'une page avec du texte natif (pas un scan)."""
    path = tmp_path / "native.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Ceci est un document de test avec du texte natif.\n" * 5)
    doc.save(str(path))
    doc.close()
    return str(path)


@pytest.fixture
def scanned_pdf_path(tmp_path):
    """PDF d'une page composée uniquement d'une image (aucun texte natif) -- simule un scan."""
    from PIL import Image, ImageDraw

    img_path = tmp_path / "scan.png"
    img = Image.new("RGB", (800, 200), color="white")
    draw = ImageDraw.Draw(img)
    draw.text((20, 80), "SCANNED TEXT FOR OCR", fill="black")
    img.save(img_path)

    path = tmp_path / "scanned.pdf"
    doc = fitz.open()
    page = doc.new_page(width=800, height=200)
    page.insert_image(page.rect, filename=str(img_path))
    doc.save(str(path))
    doc.close()
    return str(path)


@pytest.fixture
def empty_pdf_path(tmp_path):
    """PDF d'une page totalement vierge (ni texte ni image)."""
    path = tmp_path / "empty.pdf"
    doc = fitz.open()
    doc.new_page()
    doc.save(str(path))
    doc.close()
    return str(path)


@pytest.fixture
def sample_docx_path(tmp_path):
    """DOCX avec deux paragraphes et un tableau 2x2."""
    from docx import Document

    path = tmp_path / "sample.docx"
    doc = Document()
    doc.add_paragraph("Premier paragraphe de test.")
    doc.add_paragraph("Second paragraphe.")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Nom"
    table.rows[0].cells[1].text = "Valeur"
    table.rows[1].cells[0].text = "Total"
    table.rows[1].cells[1].text = "42"
    doc.save(str(path))
    return str(path)


@pytest.fixture
def empty_docx_path(tmp_path):
    """DOCX sans aucun paragraphe ni tableau."""
    from docx import Document

    path = tmp_path / "empty.docx"
    Document().save(str(path))
    return str(path)


@pytest.fixture
def sample_pptx_path(tmp_path):
    """PPTX d'une slide avec titre + contenu."""
    from pptx import Presentation

    path = tmp_path / "sample.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Titre de test"
    slide.placeholders[1].text = "Contenu de test PPTX"
    prs.save(str(path))
    return str(path)


@pytest.fixture
def sample_xlsx_path(tmp_path):
    """XLSX avec une feuille 2x2 (mêmes données que sample_docx_path)."""
    import openpyxl

    path = tmp_path / "sample.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "Nom"
    ws["B1"] = "Valeur"
    ws["A2"] = "Total"
    ws["B2"] = 42
    wb.save(str(path))
    return str(path)


@pytest.fixture
def sample_image_path(tmp_path):
    """Image PNG avec du texte à reconnaître par OCR."""
    from PIL import Image, ImageDraw

    path = tmp_path / "sample.png"
    img = Image.new("RGB", (400, 100), color="white")
    draw = ImageDraw.Draw(img)
    draw.text((10, 40), "IMAGE TEST TEXT", fill="black")
    img.save(path)
    return str(path)
